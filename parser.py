from pptx import Presentation
import re

def parser(prs):
    list = prs.slides
    dict = {}
    for el in list:
        text_slide = ''
        for shape in el.shapes:
            if not shape.has_text_frame:
                continue
            text_slide += shape.text_frame.text + ' '
        text = re.sub(r'\W', ' ', text_slide)
        dict[str(el)]=([el, text])
    return dict

prs = Presentation('1.pptx')
dict=parser(prs)
print(dict)